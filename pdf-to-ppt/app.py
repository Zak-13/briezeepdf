import os
import uuid
from flask import Flask, request, send_file
from flask_cors import CORS
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
import tempfile

app = Flask(__name__)

# Enable CORS for all domains (allow any origin to make requests)
CORS(app, resources={r"/convert": {"origins": "*", "methods": ["POST"], "allow_headers": ["Content-Type"]}})

# Use a temporary directory in a cross-platform manner
TEMP_DIR = tempfile.mkdtemp()

# The function that converts PDF to PowerPoint
def pdf_to_ppt(pdf_path, pptx_path):
    prs = Presentation()
    try:
        # Open the PDF file using PyMuPDF (fitz)
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"Error opening PDF file: {e}")
        return None

    # Iterate over each page of the PDF
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # High resolution
        img_bytes = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_bytes))

        # Generate a temporary image file for each page
        img_path = os.path.join(TEMP_DIR, f"temp_page_{page_num}.png")
        img.save(img_path)

        # Add an empty slide to the presentation
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Empty layout
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

        # Clean up temporary image file after processing the page
        os.remove(img_path)

    # Save the PowerPoint file to the temporary directory
    pptx_file = os.path.join(TEMP_DIR, pptx_path)
    try:
        prs.save(pptx_file)
    except Exception as e:
        print(f"Error saving PowerPoint file: {e}")
        return None

    # Check if the PowerPoint file exists
    if not os.path.exists(pptx_file):
        print(f"Error: PowerPoint file {pptx_file} was not saved successfully.")
        return None

    print(f"Conversion complete! PowerPoint saved as {pptx_file}")
    return pptx_file  # Return the full path to the PowerPoint file

# Route to handle file upload and conversion
@app.route('/convert', methods=['POST'])
def convert_pdf_to_ppt():
    if 'file' not in request.files:
        return "No file part", 400

    file = request.files['file']
    if file.filename == '':
        return "No selected file", 400

    # Extract the original PDF filename without extension
    pdf_filename = os.path.splitext(file.filename)[0]  # Get filename without extension

    # Generate the PDF path using the original filename
    unique_filename = str(uuid.uuid4())  # Generate a unique identifier for internal use
    pdf_path = os.path.join(TEMP_DIR, f"{pdf_filename}_{unique_filename}.pdf")
    pptx_filename = f"{pdf_filename}.pptx"  # Use the original PDF filename for the PPTX
    pptx_file_path = None  # Initialize to avoid the UnboundLocalError

    try:
        # Save the uploaded PDF temporarily
        file.save(pdf_path)

        # Convert the uploaded PDF to PowerPoint
        pptx_file_path = pdf_to_ppt(pdf_path, pptx_filename)
        if not pptx_file_path:
            return f"Error: The PowerPoint file was not saved correctly.", 500

        # Send the converted PowerPoint file back to the client with the correct filename
        return send_file(pptx_file_path, as_attachment=True, download_name=pptx_filename)

    except Exception as e:
        print(f"Error during file processing: {e}")
        return f"Error: {str(e)}", 500
    finally:
        # Clean up the uploaded PDF file
        if os.path.exists(pdf_path):
            os.remove(pdf_path)

        # Clean up the PowerPoint file if it was created
        if pptx_file_path and os.path.exists(pptx_file_path):
            os.remove(pptx_file_path)

if __name__ == '__main__':
    app.run(debug=True)
